"""PPTX extractor (python-pptx).

Each slide becomes one Page. For every shape we extract either:
  - text (title / body / placeholder) -> text block
  - table                               -> table block (Markdown)
  - picture                             -> image block for OCR
  - grouped shapes                      -> recursively flattened
Speaker notes are appended as a final text block per slide.

Legacy .ppt files are converted via LibreOffice.
"""
from __future__ import annotations

import io
import logging
from pathlib import Path
from typing import List, Optional

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .base import BaseExtractor, Block, ExtractedDoc, Page
from ..libreoffice import LibreOfficeError, convert as lo_convert
from ..scripts import detect_script

log = logging.getLogger(__name__)


class PptxExtractor(BaseExtractor):
    extensions = ["pptx", "ppt"]
    format_name = "pptx"

    def extract(self, path: str | Path) -> ExtractedDoc:
        path = Path(path)
        src = path
        if path.suffix.lower() == ".ppt":
            src = self._ppt_to_pptx(path)

        prs = Presentation(str(src))
        min_side = int(self.cfg.ocr.min_image_side_px)
        ocr_embedded = bool(self.cfg.extractors.pptx.get("ocr_embedded_images", True))

        doc = ExtractedDoc(source_path=str(path), format="pptx")

        for idx, slide in enumerate(prs.slides):
            page = Page(index=idx, label=f"Slide {idx + 1}")

            # Phase 1: gather all text on the slide (including notes +
            # grouped shapes) to detect the dominant script. Used as the
            # default script_hint for every image block on this slide when
            # detection_mode is surrounding_text / both.
            slide_text_parts: List[str] = []
            self._gather_text(slide.shapes, slide_text_parts)
            if slide.has_notes_slide:
                slide_text_parts.append(slide.notes_slide.notes_text_frame.text or "")
            slide_script = detect_script("\n".join(slide_text_parts))

            # Phase 2: build blocks, attaching script_hint to image blocks.
            self._walk_shapes(slide.shapes, page, min_side, ocr_embedded,
                              slide_script)

            # Speaker notes
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    page.blocks.append(
                        Block(type="text",
                              text=f"> **Notes:** {notes}",
                              source="pptx_notes")
                    )
            doc.pages.append(page)

        return doc

    # ------------------------------------------------------------------
    # Text-only walker (used for script detection pre-pass)
    # ------------------------------------------------------------------

    def _gather_text(self, shapes, out: List[str]) -> None:
        for shape in shapes:
            stype = shape.shape_type
            if stype == MSO_SHAPE_TYPE.GROUP:
                self._gather_text(shape.shapes, out)
                continue
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        out.append(cell.text or "")
                continue
            if shape.has_text_frame:
                out.append(shape.text_frame.text or "")

    # ------------------------------------------------------------------
    # Shape walker
    # ------------------------------------------------------------------

    def _walk_shapes(self, shapes, page: Page, min_side: int,
                     ocr_embedded: bool,
                     slide_script: Optional[str] = None) -> None:
        for shape in shapes:
            stype = shape.shape_type

            if stype == MSO_SHAPE_TYPE.GROUP:
                self._walk_shapes(shape.shapes, page, min_side,
                                  ocr_embedded, slide_script)
                continue

            if stype == MSO_SHAPE_TYPE.PICTURE:
                if not ocr_embedded:
                    continue
                img = self._load_picture(shape, min_side)
                if img is not None:
                    page.blocks.append(
                        Block(type="image", image=img, source="pptx_picture",
                              meta={"w": img.width, "h": img.height},
                              script_hint=slide_script)
                    )
                continue

            if shape.has_table:
                md = self._table_to_md(shape.table)
                if md:
                    page.blocks.append(Block(type="table", text=md,
                                             source="pptx_table"))
                continue

            if shape.has_text_frame:
                text = self._text_frame_to_md(shape.text_frame)
                if text.strip():
                    page.blocks.append(Block(type="text", text=text,
                                             source="pptx_text"))

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------

    @staticmethod
    def _text_frame_to_md(tf) -> str:
        parts: List[str] = []
        for para in tf.paragraphs:
            text = "".join(run.text for run in para.runs)
            if not text:
                continue
            level = getattr(para, "level", 0) or 0
            if level == 0:
                parts.append(text)
            else:
                parts.append(("  " * (level - 1)) + f"- {text}")
        return "\n".join(parts)

    @staticmethod
    def _table_to_md(tbl) -> str:
        rows = [[cell.text.strip().replace("\n", " ") for cell in row.cells]
                for row in tbl.rows]
        if not rows:
            return ""
        header = rows[0]
        md = ["| " + " | ".join(header) + " |",
              "| " + " | ".join(["---"] * len(header)) + " |"]
        for r in rows[1:]:
            md.append("| " + " | ".join(r) + " |")
        return "\n".join(md)

    @staticmethod
    def _load_picture(shape, min_side: int):
        try:
            blob = shape.image.blob
            img = Image.open(io.BytesIO(blob)).convert("RGB")
        except Exception as e:
            log.debug("Could not load picture: %s", e)
            return None
        if img.width < min_side or img.height < min_side:
            return None
        return img

    # ------------------------------------------------------------------
    # .ppt -> .pptx via LibreOffice
    # ------------------------------------------------------------------

    def _ppt_to_pptx(self, path: Path) -> Path:
        if not self.cfg.extractors.pptx.get("legacy_ppt_via_libreoffice", True):
            raise RuntimeError(".ppt support disabled via config")
        log.info("Converting legacy .ppt -> .pptx via LibreOffice: %s", path)
        lo_cfg = self.cfg.extractors.get("libreoffice", {}) or {}
        try:
            return lo_convert(
                path,
                target_ext="pptx",
                timeout_s=float(lo_cfg.get("timeout_s", 300)),
                retries=int(lo_cfg.get("retries", 1)),
            )
        except LibreOfficeError as e:
            raise RuntimeError(str(e))
